# -*- coding: utf-8 -*-
"""
발표용 PPTX 생성 (Word 복사 없이 report/images 직접 삽입)
출력: report/최종보고서_발표.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
REPORT = ROOT / "report"
IMG = REPORT / "images"
OUT = REPORT / "최종보고서_발표.pptx"

# (제목, 본문 불릿 리스트, 이미지 파일명 또는 None)
SLIDES = [
    (
        "서울시 차세대 트렌드 상권 발굴",
        ["분석: 423개 행정동 · 2020~2024", "20대 이동(KT) + 소비(추정매출) 양축", "Phase35·41 혼합 + Tier 분류"],
        None,
    ),
    (
        "핵심 결론",
        [
            "여가형(주말집중도≥0.8)에서만 이동→소비 전환",
            "레퍼런스 4개: 방배2·청운효자·합정·성수2가1동",
            "Tier1 후보: 숭인2동·신당5동 (교집합)",
            "임대료는 이동·소비 이후 후행 지표",
        ],
        None,
    ),
    (
        "구역 유형 · 레퍼런스",
        ["오피스형: 업종점수↑ 트렌드매출↓ → 게이트 필요", "레퍼런스 4개 = 전환 완료 기준"],
        "S2_05_area_type_industry.png",
    ),
    (
        "가설 검증: 이동·소비·lag",
        ["분기 수준 r≈0.81 (이동비율↔20대비중)", "Phase41: 이동↑ 소비↓ (초기 단계)", "성수2가1: 2021 이동 → 2022 소비 전환"],
        "S3_04_causal_lag.png",
    ),
    (
        "대표 사례: 성수2가1동",
        ["20대 이동 +236% (2020→24)", "20대 매출 +88.5%", "뚝섬 임대료 +51.1% (후행)"],
        "S3_05_ref_sales_timeseries.png",
    ),
    (
        "Phase 35 — 코사인 유사도",
        ["레퍼런스 궤적 유사 · 중기~성숙", "1순위: 숭인2동(0.944), 신당5동(0.930)"],
        "S4a_01_phase35_candidates.png",
    ),
    (
        "Phase 41 — 4분면",
        ["2사분면: 저활성 × 고이동성장", "이촌2동 이동성장세 1.74"],
        "S4b_01_quadrant.png",
    ),
    (
        "Tier 통합",
        [
            "Tier1(2): 숭인2·신당5 — 최우선",
            "Tier2(13): 업종 전환·궤적 유사",
            "Tier3(13): 이동 선행·소비 미전환",
        ],
        "S5_01_phase35_vs_41.png",
    ),
    (
        "시사점",
        [
            "단일 지표·단일 방법론 불충분",
            "주말집중도 후처리 필수",
            "모니터링: Tier별 분기 추적",
        ],
        None,
    ),
]


def add_bullets(text_frame, lines: list[str]):
    text_frame.clear()
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "맑은 고딕"


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank

    for title, bullets, img_name in SLIDES:
        slide = prs.slides.add_slide(blank)

        # 제목
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9))
        tf = box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = "맑은 고딕"

        has_img = img_name and (IMG / img_name).is_file()
        body_top = Inches(1.35)
        body_w = Inches(6.2) if has_img else Inches(12.3)

        body = slide.shapes.add_textbox(Inches(0.5), body_top, body_w, Inches(5.5))
        add_bullets(body.text_frame, bullets)

        if has_img:
            slide.shapes.add_picture(
                str(IMG / img_name),
                Inches(6.9),
                Inches(1.2),
                width=Inches(6.0),
            )
        elif img_name:
            note = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.5))
            note.text_frame.text = f"[이미지 없음: {img_name}]"

    prs.save(str(OUT))
    print(f"저장: {OUT}")


if __name__ == "__main__":
    main()
